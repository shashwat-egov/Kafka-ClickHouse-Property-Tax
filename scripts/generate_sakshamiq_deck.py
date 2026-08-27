#!/usr/bin/env python3
"""
Generate the SakshamIQ pitch deck: a branded overview of the analytical
datastore modernization product (Kafka + ClickHouse + Airflow, Bronze-
Silver-Gold medallion architecture), with Punjab Property Tax as the
exemplar implementation.

Built on eGov's own PowerPoint template so the deck inherits real eGov
branding (navy/orange color system, logo placement, fonts) rather than an
invented palette. The template's own 32 slides are removed and replaced
with new content; its masters/layouts/theme travel with the file.

Content is sourced from CLAUDE.md (implementation detail: table names,
mart list, RMV dependencies) and from the "Analytical Datastore
Modernization" design doc (problem statement, validated scale numbers,
storage/performance benchmarks, governance framework, roadmap).

All slides are native, editable PowerPoint shapes/charts/tables (not
flattened images).

Usage:
    python3 generate_sakshamiq_deck.py [--output docs/SakshamIQ.pptx]
                                        [--template ~/Downloads/Template.pptx]
"""

import argparse
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn

# Canvas - matches the eGov template's own slide size (16:9 @ 10in wide)
SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(5.625)
MARGIN = Inches(0.45)
CONTENT_RIGHT = Emu(int(SLIDE_WIDTH) - int(MARGIN))
CONTENT_TOP = Inches(1.15)
CONTENT_BOTTOM = Inches(5.32)

# eGov brand palette (extracted from Template.pptx's actual slide colors)
NAVY = RGBColor(0x27, 0x3A, 0x80)
NAVY_DARK = RGBColor(0x17, 0x2A, 0x61)
NAVY_LIGHT = RGBColor(0x6B, 0x82, 0xC2)
ORANGE = RGBColor(0xF6, 0x85, 0x21)
ORANGE_DARK = RGBColor(0xC2, 0x5E, 0x0E)
ORANGE_LIGHT = RGBColor(0xFB, 0xC1, 0x89)
TEXT_DARK = RGBColor(0x1F, 0x1F, 0x1F)
TEXT_MUTED = RGBColor(0x7F, 0x7F, 0x7F)
RULE_GREY = RGBColor(0xDD, 0xDD, 0xDD)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT_HEAD = "DM Sans"
FONT_BODY = "Lato"

ASSETS_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "docs", "assets")
IMAGE_LOGO = os.path.join(ASSETS_DIR, "egov-logo.jpg")
IMAGE_HERO = os.path.join(ASSETS_DIR, "digital-enablement-india.jpg")
LOGO_ASPECT = 3.1201413427561837
HERO_ASPECT = 1.030281182408075


# ---------------------------------------------------------------- helpers --

def delete_all_slides(prs):
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(qn("r:id"))
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def get_blank_layout(prs):
    for layout in prs.slide_layouts:
        if layout.name.strip().upper() == "BLANK":
            return layout
    return prs.slide_layouts[-1]


def new_slide():
    return _PRS.slides.add_slide(_BLANK_LAYOUT)


def set_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, size=18, bold=False,
                 italic=False, color=TEXT_DARK, align=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.TOP, font=FONT_BODY, spacing=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = font
        if spacing:
            p.line_spacing = spacing
    return box


def add_rect(slide, left, top, width, height, color, shadow=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = shadow
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color,
                      line_color=None, radius=0.08, shadow=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = radius
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    shape.shadow.inherit = shadow
    return shape


def add_ring(slide, cx, cy, diameter, color, weight=2.25):
    ring = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, diameter, diameter)
    ring.fill.background()
    ring.line.color.rgb = color
    ring.line.width = Pt(weight)
    ring.shadow.inherit = False
    return ring


def add_icon(slide, left, top, diameter, icon, bg_color, fg_color=WHITE):
    """icon: an MSO_SHAPE member (drawn inside the circle) or a short string
    (drawn as a bold monogram)."""
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diameter, diameter)
    circle.fill.solid()
    circle.fill.fore_color.rgb = bg_color
    circle.line.fill.background()
    circle.shadow.inherit = False

    if isinstance(icon, str):
        tf = circle.text_frame
        tf.word_wrap = False
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = 0
        tf.margin_right = 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = icon
        run.font.size = Pt(15 if len(icon) <= 3 else 11)
        run.font.bold = True
        run.font.color.rgb = fg_color
        run.font.name = FONT_HEAD
    else:
        inset = Emu(int(diameter * 0.3))
        inner = Emu(int(diameter) - 2 * int(inset))
        glyph = slide.shapes.add_shape(
            icon, Emu(int(left) + int(inset)), Emu(int(top) + int(inset)), inner, inner,
        )
        glyph.fill.solid()
        glyph.fill.fore_color.rgb = fg_color
        glyph.line.fill.background()
        glyph.shadow.inherit = False
    return circle


def add_logo(slide, dark_bg=False):
    w = Inches(0.95)
    h = Emu(int(w / LOGO_ASPECT))
    left = Emu(int(SLIDE_WIDTH) - int(MARGIN) - int(w))
    top = Inches(0.2)
    if dark_bg:
        pad = Inches(0.09)
        add_rounded_rect(
            slide, Emu(int(left) - int(pad)), Emu(int(top) - int(pad)),
            Emu(int(w) + 2 * int(pad)), Emu(int(h) + 2 * int(pad)),
            WHITE, radius=0.3,
        )
    slide.shapes.add_picture(IMAGE_LOGO, left, top, width=w, height=h)


def add_kicker(slide, kicker, y=Inches(0.28)):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, MARGIN, Emu(int(y) + Inches(0.03)),
                                  Inches(0.11), Inches(0.11))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ORANGE
    dot.line.fill.background()
    dot.shadow.inherit = False
    add_textbox(slide, Emu(int(MARGIN) + Inches(0.22)), y, Inches(6), Inches(0.28),
                kicker.upper(), size=10.5, bold=True, color=TEXT_MUTED, font=FONT_HEAD)


def content_slide(kicker, title, title_size=22):
    slide = new_slide()
    set_background(slide, WHITE)
    add_kicker(slide, kicker)
    add_rect(slide, MARGIN, Inches(0.56), Inches(0.07), Inches(0.42), ORANGE)
    add_textbox(slide, Emu(int(MARGIN) + Inches(0.22)), Inches(0.52), Inches(8.3), Inches(0.55),
                title, size=title_size, bold=True, color=NAVY, font=FONT_HEAD)
    add_logo(slide)
    return slide


def add_numbered_list(slide, items, left, top, width, height,
                       num_color=ORANGE, title_color=TEXT_DARK, rule=True):
    n = len(items)
    item_h = Emu(int(height) // n)
    for i, item in enumerate(items):
        title, desc = item if isinstance(item, tuple) else (item, None)
        item_top = Emu(int(top) + i * int(item_h))
        add_textbox(slide, left, item_top, Inches(0.5), Inches(0.3),
                    f"{i + 1:02d}", size=13.5, bold=True, color=num_color, font=FONT_HEAD)
        text_left = Emu(int(left) + Inches(0.48))
        text_w = Emu(int(width) - Inches(0.48))
        add_textbox(slide, text_left, item_top, text_w, Inches(0.3),
                    title, size=12.5, bold=True, color=title_color, font=FONT_HEAD)
        if desc:
            add_textbox(slide, text_left, Emu(int(item_top) + Inches(0.27)), text_w,
                        Emu(int(item_h) - Inches(0.4)),
                        desc, size=10, color=TEXT_MUTED, font=FONT_BODY)
        if rule and i < n - 1:
            rule_y = Emu(int(item_top) + int(item_h) - Inches(0.06))
            add_rect(slide, left, rule_y, width, Pt(0.75), RULE_GREY)


def add_dash_list(slide, items, left, top, width, height, size=12, color=TEXT_DARK):
    """items: each entry is either a plain string, or a (bold_prefix, rest) tuple
    when the line needs to lead with a bold code/term before the regular text."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(9)
        r1 = p.add_run()
        r1.text = "—  "
        r1.font.size = Pt(size)
        r1.font.bold = True
        r1.font.color.rgb = ORANGE
        r1.font.name = FONT_HEAD
        if isinstance(item, tuple):
            bold_prefix, rest = item
            rb = p.add_run()
            rb.text = bold_prefix
            rb.font.size = Pt(size)
            rb.font.bold = True
            rb.font.color.rgb = NAVY
            rb.font.name = FONT_HEAD
            rr = p.add_run()
            rr.text = rest
            rr.font.size = Pt(size)
            rr.font.color.rgb = color
            rr.font.name = FONT_BODY
        else:
            r2 = p.add_run()
            r2.text = item
            r2.font.size = Pt(size)
            r2.font.color.rgb = color
            r2.font.name = FONT_BODY
    return box


def add_callout_bar(slide, left, top, width, height, label, text,
                     bg_color=RGBColor(0xEE, 0xF1, 0xFA), label_color=NAVY,
                     text_color=TEXT_DARK, text_bold=False):
    add_rounded_rect(slide, left, top, width, height, bg_color, radius=0.2)
    box = slide.shapes.add_textbox(Emu(int(left) + Inches(0.22)), top,
                                    Emu(int(width) - Inches(0.4)), height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    if label:
        r1 = p.add_run()
        r1.text = label.upper() + "   "
        r1.font.bold = True
        r1.font.size = Pt(10.5)
        r1.font.color.rgb = label_color
        r1.font.name = FONT_HEAD
    r2 = p.add_run()
    r2.text = text
    r2.font.bold = text_bold
    r2.font.size = Pt(10.5)
    r2.font.color.rgb = text_color
    r2.font.name = FONT_BODY
    return box


def add_icon_panel(slide, icon, color, cx, cy, diameter=Inches(1.7)):
    left = Emu(int(cx) - int(diameter) // 2)
    top = Emu(int(cy) - int(diameter) // 2)
    add_icon(slide, left, top, diameter, icon, color)


def add_stat_tiles(slide, stats, top, height=Inches(1.9)):
    n = len(stats)
    gap = Inches(0.25)
    avail_w = Emu(int(CONTENT_RIGHT) - int(MARGIN))
    tile_w = Emu((int(avail_w) - (n - 1) * int(gap)) // n)

    for i, (number, label) in enumerate(stats):
        left = Emu(int(MARGIN) + i * (int(tile_w) + int(gap)))
        add_rounded_rect(slide, left, top, tile_w, height, WHITE, RULE_GREY, radius=0.1)
        add_textbox(slide, left, Emu(int(top) + Inches(0.22)), tile_w, Inches(0.75),
                    number, size=27, bold=True, color=ORANGE, font=FONT_HEAD,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Emu(int(left) + Inches(0.1)), Emu(int(top) + Inches(1.0)),
                    Emu(int(tile_w) - Inches(0.2)), Inches(0.75),
                    label, size=10, color=TEXT_MUTED, align=PP_ALIGN.CENTER, font=FONT_BODY)


def add_table(slide, headers, rows, left, top, width, row_h=0.42, col_widths=None):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    height = Inches(row_h * n_rows)
    gtable = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    if col_widths:
        for i, w in enumerate(col_widths):
            gtable.columns[i].width = Inches(w)

    for c, header in enumerate(headers):
        cell = gtable.cell(0, c)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.margin_top = Pt(4)
        cell.margin_bottom = Pt(4)
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(11)
                r.font.color.rgb = WHITE
                r.font.name = FONT_HEAD

    for r_idx, row in enumerate(rows, start=1):
        for c_idx, value in enumerate(row):
            cell = gtable.cell(r_idx, c_idx)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xF7, 0xF7, 0xF9) if r_idx % 2 == 0 else WHITE
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(10.5)
                    r.font.color.rgb = TEXT_DARK
                    r.font.name = FONT_BODY
    return gtable


# ------------------------------------------------------------- slide sets --

def add_title_slide():
    slide = new_slide()
    set_background(slide, WHITE)
    add_logo(slide)

    add_textbox(slide, Inches(0.55), Inches(0.95), Inches(6.2), Inches(0.95),
                "SakshamIQ", size=40, bold=True, color=NAVY, font=FONT_HEAD)
    add_textbox(slide, Inches(0.55), Inches(1.7), Inches(6.2), Inches(0.5),
                "( Analytical Intelligence for Digital Governance )",
                size=15, bold=True, color=ORANGE, font=FONT_HEAD)
    add_textbox(slide, Inches(0.55), Inches(2.3), Inches(6.0), Inches(0.9),
                "Built on Apache Kafka, ClickHouse & Apache Airflow for real-time, "
                "governed analytics at state scale.",
                size=12.5, color=TEXT_DARK, font=FONT_BODY)
    add_textbox(slide, Inches(0.55), Inches(4.75), Inches(6.5), Inches(0.4),
                "EXEMPLAR: PUNJAB PROPERTY TAX", size=12, bold=True, color=ORANGE,
                font=FONT_HEAD)

    hero_w = Inches(3.15)
    hero_h = Emu(int(hero_w / HERO_ASPECT))
    slide.shapes.add_picture(IMAGE_HERO, Emu(int(SLIDE_WIDTH) - hero_w - Inches(0.2)),
                              Emu(int(SLIDE_HEIGHT) - int(hero_h)), width=hero_w, height=hero_h)


def add_section_divider(kicker, title):
    slide = new_slide()
    set_background(slide, NAVY)
    add_ring(slide, Emu(int(SLIDE_WIDTH) - Inches(1.7)), Emu(int(SLIDE_HEIGHT) - Inches(1.7)),
             Inches(2.6), ORANGE, weight=2)
    add_ring(slide, Emu(int(SLIDE_WIDTH) - Inches(1.15)), Emu(int(SLIDE_HEIGHT) - Inches(1.15)),
             Inches(1.5), NAVY_LIGHT, weight=2)
    add_ring(slide, Inches(-0.9), Inches(-0.9), Inches(1.8), NAVY_LIGHT, weight=2)

    add_textbox(slide, Inches(0.8), Inches(2.0), Inches(8.4), Inches(0.35),
                kicker.upper(), size=13, bold=True, color=ORANGE, font=FONT_HEAD,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.8), Inches(2.4), Inches(8.4), Inches(0.9),
                title, size=32, bold=True, color=WHITE, font=FONT_HEAD,
                align=PP_ALIGN.CENTER)
    add_logo(slide, dark_bg=True)


def add_problem_slide():
    slide = content_slide("The Problem", "A Decade of Growth Outgrew the Reporting Stack")
    items = [
        ("No Cross-Domain Joins", "The existing analytical store cannot join across "
         "domains, constraining cross-cutting analysis."),
        ("Falling Back to OLTP", "State teams routinely extract reports directly from "
         "the production PostgreSQL database."),
        ("Explosive Data Growth", "A decade of demand, collection, and property data has "
         "made ad hoc reporting increasingly expensive."),
        ("Join-Heavy Queries Stall", "Reporting queries against the transactional schema "
         "have become progressively slower."),
        ("Production Under Pressure", "Analytical workloads compete with live transactions "
         "for CPU and memory on the OLTP system."),
    ]
    add_numbered_list(slide, items, MARGIN, CONTENT_TOP, Inches(5.9), Inches(3.85))
    add_icon_panel(slide, MSO_SHAPE.NO_SYMBOL, ORANGE_DARK,
                   Inches(8.05), Inches(3.0), Inches(1.6))
    add_textbox(slide, MARGIN, Inches(5.02), Inches(9.1), Inches(0.3),
                "Triggered most visibly during the One-Time Settlement (OTS) drive.",
                size=10.5, italic=True, color=TEXT_MUTED)


def add_why_punjab_slide():
    slide = content_slide("Why Punjab", "The Exemplar for Statewide Modernization")
    add_stat_tiles(slide, [
        ("2013", "Property Tax module live since"),
        ("10+ yrs", "of continuous demand, collection\n& property data"),
        ("1st", "DIGIT implementation chosen to\nvalidate this approach"),
    ], top=Inches(1.5), height=Inches(1.85))
    add_textbox(slide, MARGIN, Inches(3.65), Inches(9.1), Inches(1.4),
                "Punjab's Property Tax system carries over a decade of operational "
                "history — enough real volume and complexity to prove the analytical "
                "architecture before it is rolled out to every DIGIT urban module.",
                size=13, color=TEXT_DARK)


def add_requirements_slide():
    slide = content_slide("What We're Building", "What SakshamIQ Must Deliver", title_size=20)
    left_items = [
        ("Separate OLTP / OLAP", "Isolate analytics from transactional systems."),
        ("Event-Driven, Append-Only", "Every change flows in as an immutable event."),
        ("SQL-Native Access", "Broad accessibility for analysts via standard SQL."),
        ("Efficient Aggregation", "Fast, reliable KPI extraction at state scale."),
    ]
    right_items = [
        ("Governance Built In", "RBAC, row-level security, full auditability."),
        ("Safe Schema Evolution", "Structural change without breaking consumers."),
        ("Long-Term Scalability", "Grows with a decade of statewide data."),
        ("Proportionate Cost", "Scale without disproportionate infra cost."),
    ]
    col_w = Inches(4.35)
    add_numbered_list(slide, left_items, MARGIN, CONTENT_TOP, col_w, Inches(3.9))
    add_numbered_list(slide, right_items, Emu(int(MARGIN) + Inches(4.75)), CONTENT_TOP,
                       col_w, Inches(3.9))


def add_why_clickhouse_slide():
    slide = content_slide("Technology Choice", "Why ClickHouse?", title_size=22)
    add_textbox(slide, MARGIN, Inches(1.0), Inches(9.1), Inches(0.4),
                "The analytical workload needs a dedicated OLAP store — separate from "
                "transactional PostgreSQL.", size=11.5, color=TEXT_DARK)
    items = [
        ("Separate OLTP / OLAP", "Isolate analytical workloads from the transactional "
         "system so reporting no longer competes with live transactions."),
        ("SQL-based Analytics", "Provide SQL access for analysts and enable broader "
         "analytical accessibility."),
        ("Efficient KPI Extraction", "Support large-scale aggregations and precomputed "
         "KPI refreshes for reporting and dashboards."),
        ("Scale + Cost", "Support long-term growth without disproportionate infrastructure "
         "cost; validation was performed at ~10-15× production scale."),
    ]
    add_numbered_list(slide, items, MARGIN, Inches(1.5), Inches(5.9), Inches(3.0))
    add_icon_panel(slide, MSO_SHAPE.CAN, NAVY, Inches(8.05), Inches(3.0), Inches(1.6))
    add_callout_bar(slide, MARGIN, Inches(4.75), Inches(9.1), Inches(0.45),
                     "Validated at Scale",
                     "~50M properties · ~724M PT demands · ~10 years historical accumulation")


def add_why_stack_slide():
    slide = content_slide("Technology Stack", "Why Kafka + ClickHouse + Airflow?",
                           title_size=20)
    add_textbox(slide, MARGIN, Inches(1.0), Inches(9.1), Inches(0.4),
                "Each component has one clear responsibility in the analytical platform.",
                size=11.5, color=TEXT_DARK)

    roles = [
        ("01", "DIGIT Services", "Generates business transactions and upstream "
         "create/update events.", NAVY_DARK),
        ("02", "Apache Kafka", "Event backbone that transports immutable upstream "
         "changes.", NAVY),
        ("03", "ClickHouse", "Analytical datastore: Bronze → Silver → Gold business "
         "marts.", NAVY_DARK),
        ("04", "Apache Airflow", "Control plane: scheduling, dependencies, retries, "
         "backfills.", NAVY),
        ("05", "Dashboards / APIs", "Consume Gold marts and query Silver for analytical "
         "exploration.", NAVY_DARK),
    ]
    box_w = Inches(1.68)
    gap = Inches(0.13)
    box_h = Inches(2.15)
    top = Inches(1.5)
    left = Inches(0.45)
    for num, title, desc, fill in roles:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, box_w, box_h)
        box.adjustments[0] = 0.08
        box.fill.solid()
        box.fill.fore_color.rgb = fill
        box.line.fill.background()
        box.shadow.inherit = False
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.12)
        tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.14)
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = num
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = ORANGE
        r.font.name = FONT_HEAD
        p2 = tf.add_paragraph()
        p2.space_before = Pt(4)
        r2 = p2.add_run()
        r2.text = title
        r2.font.size = Pt(12)
        r2.font.bold = True
        r2.font.color.rgb = WHITE
        r2.font.name = FONT_HEAD
        p3 = tf.add_paragraph()
        p3.space_before = Pt(6)
        r3 = p3.add_run()
        r3.text = desc
        r3.font.size = Pt(8.5)
        r3.font.color.rgb = RGBColor(0xE3, 0xE7, 0xF4)
        r3.font.name = FONT_BODY
        left = Emu(int(left) + int(box_w) + int(gap))

    add_callout_bar(slide, MARGIN, Inches(3.95), Inches(9.1), Inches(0.55),
                     "Control plane, not data store",
                     "Airflow orchestrates execution; ClickHouse persists and processes "
                     "analytical data.", bg_color=RGBColor(0xFC, 0xEC, 0xDA),
                     label_color=ORANGE_DARK)


def add_principles_slide():
    slide = content_slide("The Architecture", "Four Design Principles")
    items = [
        ("Event-Driven & Append-Only", "All creates/updates flow as immutable events — "
         "no updates or deletes, full historical traceability."),
        ("Layered (Medallion) Model", "Raw Events → Canonical Entities → Data Marts, each "
         "layer with one job and one guarantee."),
        ("Governance-First Design", "RBAC, row-level isolation, snapshot immutability, and "
         "a versioned mart strategy from day one."),
        ("Operational Simplicity", "Airflow is reused for orchestration; ClickHouse is the "
         "only new stateful component."),
    ]
    add_numbered_list(slide, items, MARGIN, CONTENT_TOP, Inches(5.9), Inches(3.85))
    add_icon_panel(slide, MSO_SHAPE.HEXAGON, NAVY, Inches(8.05), Inches(3.0), Inches(1.6))


def add_layer_box(slide, left, top, width, height, label, sublabel, fill,
                   text_color=WHITE, outline=None, label_size=10.5, sublabel_size=7.5):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.adjustments[0] = 0.1
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    if outline is None:
        box.line.fill.background()
    else:
        box.line.color.rgb = outline
        box.line.width = Pt(1.5)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(label_size)
    run.font.bold = True
    run.font.color.rgb = text_color
    run.font.name = FONT_HEAD
    if sublabel:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = sublabel
        run2.font.size = Pt(sublabel_size)
        run2.font.color.rgb = text_color
        run2.font.name = FONT_BODY


def add_flow_arrow(slide, left, top, width, height, color=TEXT_MUTED):
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()
    arrow.shadow.inherit = False
    return arrow


def add_diagram_slide():
    slide = content_slide("The Architecture", "End-to-End Data Flow", title_size=20)

    box_top = Inches(2.15)
    box_h = Inches(1.15)
    box_w = Inches(1.32)
    gap = Inches(0.16)

    steps = [
        ("DIGIT Services", "Property & demand\nevents (via Kafka)", NAVY_DARK),
        ("Ingestion", "Streamed, no parsing", NAVY),
        ("Raw Events", "BRONZE · append-only", ORANGE_DARK),
        ("Canonical Entities", "SILVER · deduplicated", NAVY_LIGHT),
        ("Business Marts", "GOLD · pre-aggregated", ORANGE),
        ("Dashboards & APIs", "Business analysts", NAVY_DARK),
    ]

    left = Inches(0.4)
    positions = []
    for label, sublabel, fill in steps:
        add_layer_box(slide, left, box_top, box_w, box_h, label, sublabel, fill)
        positions.append(left)
        left = Emu(int(left) + int(box_w) + int(gap))

    arrow_h = Inches(0.24)
    arrow_top = Emu(int(box_top) + int(box_h) // 2 - int(arrow_h) // 2)
    for i in range(len(steps) - 1):
        a_left = Emu(int(positions[i]) + int(box_w))
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, a_left, arrow_top, gap, arrow_h)
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = TEXT_MUTED
        arrow.line.fill.background()
        arrow.shadow.inherit = False

    add_textbox(slide, Emu(int(positions[2]) - Inches(0.35)), Inches(1.15),
                Inches(2.2), Inches(0.9),
                "Airflow (T-1 nightly)\nparses raw JSON, upserts\ncanonical entities",
                size=9, color=ORANGE_DARK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Emu(int(positions[3]) + Inches(0.35)), Inches(1.15),
                Inches(2.4), Inches(0.9),
                "Dependency-ordered mart\nrefresh — consistent\nsnapshot every time",
                size=9, color=NAVY, align=PP_ALIGN.CENTER)

    add_textbox(slide, MARGIN, Inches(3.65), Inches(9.1), Inches(1.4),
                "Airflow acts strictly as a control plane — it orchestrates, but never "
                "persists analytical data. Every stage is idempotent and independently "
                "re-runnable, so a bad run never means starting over.",
                size=11.5, color=TEXT_DARK)


def add_layer_detail_slide(title, accent, icon, points):
    slide = content_slide("The Architecture", title, title_size=20)
    add_numbered_list(slide, points, MARGIN, CONTENT_TOP, Inches(5.9), Inches(3.85),
                       num_color=accent)
    add_icon_panel(slide, icon, accent, Inches(8.05), Inches(3.0), Inches(1.6))


def add_billing_events_problem_slide():
    slide = content_slide("Pipeline Enablement",
                           "Enable Billing & Collection Events for ClickHouse",
                           title_size=18)
    left_w = Inches(4.35)
    left_x = MARGIN
    right_x = Emu(int(MARGIN) + int(left_w) + Inches(0.3))

    add_textbox(slide, left_x, CONTENT_TOP, left_w, Inches(0.35),
                "Problem", size=13.5, bold=True, color=ORANGE_DARK, font=FONT_HEAD)
    add_dash_list(slide, [
        "Demand and bill data was persisted to PostgreSQL through JDBC, but was never "
        "emitted to Kafka for analytics",
        "As a result, these records were not available in the ClickHouse pipeline",
        "Existing application Kafka topics could not be reused because they serve "
        "application / integration flows",
    ], left_x, Emu(int(CONTENT_TOP) + Inches(0.4)), left_w, Inches(3.5), size=10.5)

    add_textbox(slide, right_x, CONTENT_TOP, left_w, Inches(0.35),
                "Changes Required", size=13.5, bold=True, color=NAVY, font=FONT_HEAD)
    add_dash_list(slide, [
        ("save-demand-event", " → demand creation"),
        ("update-demand-event", " → demand updates, including payment-driven updates "
         "and reversals"),
        ("save-bill-event", " → collection-side bill creation"),
        "Events are published only after successful DB persistence",
        "Existing persistence and existing Kafka flows remain unchanged",
    ], right_x, Emu(int(CONTENT_TOP) + Inches(0.4)), left_w, Inches(3.9), size=10.5)

    add_rect(slide, Emu(int(MARGIN) + int(left_w) + Inches(0.12)), CONTENT_TOP,
             Pt(1), Inches(3.9), RULE_GREY)


def add_billing_code_changes_slide():
    slide = content_slide("Code Changes", "Code Changes in Billing & Collection Services",
                           title_size=18)
    add_table(
        slide,
        ["Service", "Code Change", "Purpose"],
        [
            ["Billing", "application.properties", "Added demand event topic configuration"],
            ["", "ApplicationProperties.java", "Added topic properties / getters"],
            ["", "DemandService.create()", "Publish save-demand-event after DB save"],
            ["", "DemandService.update()", "Publish update-demand-event after DB update"],
            ["Collection", "application.properties", "Added bill event topic configuration"],
            ["", "ApplicationProperties.java", "Added bill topic property / getter"],
            ["", "PaymentService.createPayment()", "Publish save-bill-event after "
             "egcl_bill persistence"],
        ],
        left=MARGIN, top=Inches(1.3), width=Inches(9.1), row_h=0.36,
        col_widths=[1.5, 3.0, 4.6],
    )

    add_textbox(slide, MARGIN, Inches(4.2), Inches(9.1), Inches(0.3),
                "Key Design Decisions", size=12, bold=True, color=ORANGE_DARK,
                font=FONT_HEAD)
    col_w = Inches(4.35)
    add_dash_list(slide, [
        ("DemandService.update()", " is the single publishing point for all demand updates"),
        "No changes to existing save-demand, update-demand or payment topics",
    ], MARGIN, Inches(4.5), col_w, Inches(0.8), size=9.5)
    add_dash_list(slide, [
        ("save-bill-event", " is emitted from Collection Service — egcl_bill is the bill "
         "used in the payment flow"),
        "No separate bill-cancel-event / update-bill-event required",
    ], Emu(int(MARGIN) + int(col_w) + Inches(0.3)), Inches(4.5), col_w, Inches(0.8), size=9.5)


def add_billing_flow_slide():
    slide = content_slide("Resulting Flow", "How Billing & Collection Events Reach ClickHouse",
                           title_size=17)

    row_h = Inches(0.72)
    row_gap = Inches(0.1)
    top1 = Inches(1.45)
    top2 = Emu(int(top1) + int(row_h) + int(row_gap))
    top3 = Emu(int(top2) + int(row_h) + int(row_gap))
    tall_h = Emu(int(top3) + int(row_h) - int(top1))

    col_gap = Inches(0.13)
    w1, w2, w3, w4, w5 = Inches(1.35), Inches(1.55), Inches(1.6), Inches(1.35), Inches(1.35)
    x1 = Inches(0.5)
    x2 = Emu(int(x1) + int(w1) + int(col_gap))
    x3 = Emu(int(x2) + int(w2) + int(col_gap))
    x4 = Emu(int(x3) + int(w3) + int(col_gap))
    x5 = Emu(int(x4) + int(w4) + int(col_gap))

    add_layer_box(slide, x1, top1, w1, Emu(int(top2) + int(row_h) - int(top1)),
                  "Billing Service", "demand lifecycle", NAVY_DARK, label_size=9.5, sublabel_size=7)
    add_layer_box(slide, x1, top3, w1, row_h, "Collection Service", "payment / billing",
                  NAVY_DARK, label_size=9.5, sublabel_size=7)

    add_layer_box(slide, x2, top1, w2, row_h, "Demand Create", "DemandService.create()",
                  NAVY, label_size=9, sublabel_size=6.5)
    add_layer_box(slide, x2, top2, w2, row_h, "Demand Update", "DemandService.update()",
                  NAVY, label_size=9, sublabel_size=6.5)
    add_layer_box(slide, x2, top3, w2, row_h, "createPayment()", "persists egcl_bill",
                  NAVY, label_size=9, sublabel_size=6.5)

    add_layer_box(slide, x3, top1, w3, row_h, "save-demand-event", "new demand",
                  ORANGE_DARK, label_size=9, sublabel_size=6.5)
    add_layer_box(slide, x3, top2, w3, row_h, "update-demand-event",
                  "incl. payment-driven reversals", ORANGE_DARK, label_size=9, sublabel_size=6.5)
    add_layer_box(slide, x3, top3, w3, row_h, "save-bill-event", "collection-side bill",
                  ORANGE_DARK, label_size=9, sublabel_size=6.5)

    add_layer_box(slide, x4, top1, w4, tall_h, "Kafka", "dedicated analytics topics", NAVY,
                  label_size=10.5, sublabel_size=7.5)
    add_layer_box(slide, x5, top1, w5, tall_h, "ClickHouse", "Bronze → Silver → Gold",
                  ORANGE, label_size=10.5, sublabel_size=7.5)

    arrow_h = Inches(0.2)
    for row_top in (top1, top2, top3):
        ay = Emu(int(row_top) + int(row_h) // 2 - int(arrow_h) // 2)
        add_flow_arrow(slide, Emu(int(x1) + int(w1)), ay, col_gap, arrow_h)
        add_flow_arrow(slide, Emu(int(x2) + int(w2)), ay, col_gap, arrow_h)
        add_flow_arrow(slide, Emu(int(x3) + int(w3)), ay, col_gap, arrow_h)
    mid_ay = Emu(int(top1) + int(tall_h) // 2 - int(arrow_h) // 2)
    add_flow_arrow(slide, Emu(int(x4) + int(w4)), mid_ay, col_gap, arrow_h)

    outcome_top = Emu(int(top3) + int(row_h) + Inches(0.25))
    box = slide.shapes.add_textbox(MARGIN, outcome_top, Inches(9.1), Inches(0.6))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "OUTCOME   "
    r1.font.bold = True
    r1.font.size = Pt(10.5)
    r1.font.color.rgb = ORANGE_DARK
    r1.font.name = FONT_HEAD
    r2 = p.add_run()
    r2.text = ("Demand, bill and payment-related state changes are now available to the "
               "ClickHouse analytics pipeline — without modifying existing transactional "
               "persistence flows.")
    r2.font.size = Pt(10.5)
    r2.font.color.rgb = TEXT_DARK
    r2.font.name = FONT_BODY


def add_orchestration_slide():
    slide = content_slide("The Architecture", "Orchestration — Airflow as Control Plane",
                           title_size=19)
    items = [
        ("Stage-Wise Pipeline", "Raw → Entity → Mart, executed as ordered stages every run."),
        ("Scheduled Daily Execution", "T-1 batch processing runs nightly, off-peak."),
        ("Dependency-Enforced Sequencing", "Marts that depend on others always refresh in order."),
        ("Idempotent Recomputation", "Any stage can be safely re-run without side effects."),
        ("Backfill Support", "Historical reprocessing is a first-class operation."),
        ("Retry & Observability", "Automatic retry on failure, full visibility in the Airflow UI."),
    ]
    add_numbered_list(slide, items, MARGIN, CONTENT_TOP, Inches(5.9), Inches(3.85))
    add_icon_panel(slide, MSO_SHAPE.GEAR_6, NAVY, Inches(8.05), Inches(3.0), Inches(1.6))


def add_dag_bronze_to_silver_slide():
    slide = content_slide("Orchestration · DAG 1 of 2",
                           "Bronze to Silver — Five Branches, One Barrier", title_size=17)

    streams = [
        ("Property", "4 entity tables"),
        ("Demand", "1 entity table"),
        ("Payment", "1 entity table"),
        ("Bill", "2 entity tables"),
        ("Assessment", "1 entity table"),
    ]
    n = len(streams)
    row_h = Inches(0.42)
    row_gap = Inches(0.07)
    top0 = Inches(1.5)
    tall_h = Emu(n * int(row_h) + (n - 1) * int(row_gap))

    col_gap = Inches(0.12)
    w0, wA, wB, wC, wD = Inches(1.25), Inches(1.1), Inches(1.25), Inches(1.55), Inches(1.25)
    x0 = Inches(0.4)
    xA = Emu(int(x0) + int(w0) + int(col_gap))
    xB = Emu(int(xA) + int(wA) + int(col_gap))
    xC = Emu(int(xB) + int(wB) + int(col_gap))
    xD = Emu(int(xC) + int(wC) + int(col_gap))

    add_layer_box(slide, x0, top0, w0, tall_h, "Nightly Run", "T-1 window", NAVY_DARK,
                  label_size=10, sublabel_size=7.5)
    add_layer_box(slide, xD, top0, wD, tall_h, "Join & Trigger", "all five must succeed",
                  ORANGE, label_size=10, sublabel_size=7.5)

    arrow_h = Inches(0.18)
    for i, (stream, entity_count) in enumerate(streams):
        row_top = Emu(int(top0) + i * (int(row_h) + int(row_gap)))
        add_layer_box(slide, xA, row_top, wA, row_h, stream, "", RGBColor(0xEA, 0xEC, 0xF2),
                      text_color=TEXT_DARK, label_size=9.5)
        add_layer_box(slide, xB, row_top, wB, row_h, "Extract", "count the window", NAVY,
                      label_size=8.5, sublabel_size=6.5)
        add_layer_box(slide, xC, row_top, wC, row_h, "Transform & Load",
                      f"→ {entity_count}", NAVY, outline=ORANGE, label_size=8, sublabel_size=6.5)
        ay = Emu(int(row_top) + int(row_h) // 2 - int(arrow_h) // 2)
        add_flow_arrow(slide, Emu(int(x0) + int(w0)), ay, col_gap, arrow_h)
        add_flow_arrow(slide, Emu(int(xA) + int(wA)), ay, col_gap, arrow_h)
        add_flow_arrow(slide, Emu(int(xB) + int(wB)), ay, col_gap, arrow_h)
        add_flow_arrow(slide, Emu(int(xC) + int(wC)), ay, col_gap, arrow_h)

    add_textbox(slide, MARGIN, Emu(int(top0) + int(tall_h) + Inches(0.25)),
                Inches(9.1), Inches(0.8),
                "Every branch is independent — a slow stream never blocks the others, and "
                "a shared Airflow pool caps how many transform steps run at once. The gold "
                "refresh is triggered only after all five branches succeed, and this DAG "
                "completes without waiting for it.", size=10.5, color=TEXT_DARK)


def add_dag_silver_to_gold_slide():
    slide = content_slide("Orchestration · DAG 2 of 2",
                           "Silver to Gold — 13 Refreshes at Once, Then 3 That Wait",
                           title_size=16)

    rows = [
        ("2", "Property marts", None),
        ("3", "Demand & collection marts", None),
        ("1", "Demand coverage base mart", ("1", "mart-on-mart view", "waits for its base")),
        ("1", "Property change-metrics base mart",
         ("2", "mart-on-mart views", "wait for their base")),
        ("6", "Assessment, payment & rebate marts", None),
    ]
    n = len(rows)
    row_h = Inches(0.4)
    row_gap = Inches(0.08)
    top0 = Inches(1.45)
    tall_h = Emu(n * int(row_h) + (n - 1) * int(row_gap))

    w0, w1, w2 = Inches(1.2), Inches(4.3), Inches(2.5)
    gap = Inches(0.15)
    x0 = Inches(0.4)
    x1 = Emu(int(x0) + int(w0) + int(gap))
    x2 = Emu(int(x1) + int(w1) + int(gap))

    add_layer_box(slide, x0, top0, w0, tall_h, "Triggered", "when silver lands", NAVY_DARK,
                  label_size=10, sublabel_size=7.5)

    arrow_h = Inches(0.16)
    for i, (count, label, dependent) in enumerate(rows):
        row_top = Emu(int(top0) + i * (int(row_h) + int(row_gap)))
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x1, row_top, w1, row_h)
        box.adjustments[0] = 0.15
        box.fill.solid()
        box.fill.fore_color.rgb = NAVY
        box.line.fill.background()
        box.shadow.inherit = False
        tf = box.text_frame
        tf.margin_left = Inches(0.15)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        r1_ = p.add_run()
        r1_.text = count + "   "
        r1_.font.bold = True
        r1_.font.size = Pt(12)
        r1_.font.color.rgb = ORANGE
        r1_.font.name = FONT_HEAD
        r2_ = p.add_run()
        r2_.text = label
        r2_.font.bold = True
        r2_.font.size = Pt(10.5)
        r2_.font.color.rgb = WHITE
        r2_.font.name = FONT_HEAD

        if dependent:
            ay = Emu(int(row_top) + int(row_h) // 2 - int(arrow_h) // 2)
            add_flow_arrow(slide, Emu(int(x1) + int(w1)), ay, gap, arrow_h, color=ORANGE)
            dep_count, dep_label, dep_sub = dependent
            add_layer_box(slide, x2, row_top, w2, row_h, f"{dep_count} {dep_label}",
                          dep_sub, ORANGE, label_size=9.5, sublabel_size=7.5)

    end_top = Emu(int(top0) + int(tall_h) + Inches(0.18))
    add_callout_bar(slide, MARGIN, end_top, Inches(9.1), Inches(0.4), "end",
                     "a failed mart is isolated — the rest of the wave still lands",
                     bg_color=RGBColor(0x8A, 0x90, 0x9C), label_color=WHITE, text_color=WHITE)

    add_textbox(slide, MARGIN, Emu(int(end_top) + Inches(0.5)), Inches(9.1), Inches(0.7),
                "Each refresh is one SQL statement plus a watchdog: Airflow issues it, "
                "then polls ClickHouse until the view reports complete. The marts are "
                "declared manual-refresh, so Airflow is the only scheduler.",
                size=10, color=TEXT_DARK)


def add_scale_slide():
    slide = content_slide("Validated at Scale", "Proven at 10-15× Current Production Volume",
                           title_size=19)
    add_stat_tiles(slide, [
        ("5 Cr", "properties validated (~50M)"),
        ("72 Cr", "PT demands validated (~724M)"),
        ("10 yrs", "of historical accumulation"),
        ("10-15×", "current production scale"),
    ], top=Inches(1.5), height=Inches(1.85))
    add_textbox(slide, MARGIN, Inches(3.65), Inches(9.1), Inches(1.4),
                "Today's production baseline is ~33 lakh properties and ~5 crore demands. "
                "The validation dataset was built at 10-15× that volume to prove SakshamIQ "
                "holds up at full statewide, multi-year scale.",
                size=12.5, color=TEXT_DARK)


def add_storage_slide():
    slide = content_slide("Validated at Scale", "Storage Efficiency")

    chart_data = CategoryChartData()
    chart_data.categories = ["PostgreSQL\n(current scale)", "ClickHouse\n(10-15× that scale)"]
    chart_data.add_series("Storage (GB)", (284, 150))

    x, y, cx, cy = Inches(0.4), Inches(1.3), Inches(5.4), Inches(3.75)
    gframe = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)
    chart = gframe.chart
    chart.has_legend = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.number_format = '0" GB"'
    plot.data_labels.number_format_is_linked = False
    plot.data_labels.font.size = Pt(12)
    plot.data_labels.font.bold = True
    plot.data_labels.font.color.rgb = NAVY
    series = plot.series[0]
    point_colors = [NAVY_LIGHT, ORANGE]
    for i, point in enumerate(series.points):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = point_colors[i]
    chart.category_axis.tick_labels.font.size = Pt(9.5)
    chart.value_axis.tick_labels.font.size = Pt(9)

    card_left = Inches(6.05)
    add_rounded_rect(slide, card_left, Inches(1.3), Inches(3.5), Inches(3.75), WHITE, RULE_GREY)
    add_textbox(slide, Emu(int(card_left) + Inches(0.25)), Inches(1.5),
                Inches(3.0), Inches(0.4), "The Takeaway", size=13, bold=True, color=ORANGE,
                font=FONT_HEAD)
    add_dash_list(slide, [
        "ClickHouse holds ~10-15× more data than PostgreSQL's current footprint",
        "...while using ~150 GB against PostgreSQL's ~284 GB",
        "Columnar storage and compression drive the difference",
        "150 GB used of 500 GB provisioned",
    ], Emu(int(card_left) + Inches(0.25)), Inches(1.95), Inches(3.0), Inches(3.0), size=10.5)


def add_kpi_performance_slide():
    slide = content_slide("Validated at Scale", "Query Performance at Full Scale",
                           title_size=19)
    add_table(
        slide,
        ["KPI", "Dataset Size", "Refresh", "Peak Mem"],
        [
            ["Active Property Distribution Summary", "~50M properties", "~5 sec", "~142 MB"],
            ["New Property Additions by FY", "~50M properties", "~1 sec", "~132 MB"],
            ["Financial Demand & Collection (FY)", "~724M demands", "~166 sec", "~178 MB"],
            ["Monthly Collection Trend Summary", "~724M demands", "~87 sec", "~132 MB"],
            ["Property Demand Coverage Ratio (FY)", "~724M demands", "~15 sec", "~93 MB"],
            ["Property Tax Defaulters Register", "~724M demands", "~267 sec", "~196 MB"],
        ],
        left=MARGIN, top=Inches(1.35), width=Inches(9.1), row_h=0.44,
        col_widths=[4.4, 1.9, 1.4, 1.4],
    )
    add_textbox(slide, MARGIN, Inches(4.65), Inches(9.1), Inches(0.6),
                "These aggregations run only during the nightly refresh window. Every "
                "dashboard query after that reads an already-computed, lightweight mart.",
                size=10.5, italic=True, color=TEXT_MUTED)


def add_validated_airflow_slide():
    slide = content_slide("Validation", "Validated Airflow Pipeline at Scale", title_size=19)
    add_textbox(slide, MARGIN, Inches(1.0), Inches(9.1), Inches(0.4),
                "End-to-end orchestration validation across Property, Demand, Bill, "
                "Payment and Assessment streams.", size=11.5, color=TEXT_DARK)

    stats = [
        ("1M events per stream", "Property · Demand · Bill · Payment · Assessment"),
        ("18 min 23 sec", "End-to-end Raw → Silver → Gold pipeline execution"),
        ("≤ 429 MiB peak worker memory", "Well below the 1 GiB request / 2 GiB limit"),
        ("85% peak worker CPU", "Highest observed utilization during Property transformation"),
    ]
    block_w = Inches(4.3)
    block_h = Inches(0.9)
    col_gap = Inches(0.2)
    row_gap = Inches(0.15)
    top0 = Inches(1.55)
    for i, (heading, desc) in enumerate(stats):
        r, c = divmod(i, 2)
        bx = Emu(int(MARGIN) + c * (int(block_w) + int(col_gap)))
        by = Emu(int(top0) + r * (int(block_h) + int(row_gap)))
        add_rect(slide, bx, by, Inches(0.06), block_h, ORANGE)
        add_textbox(slide, Emu(int(bx) + Inches(0.2)), by, Emu(int(block_w) - Inches(0.2)),
                    Inches(0.45), heading, size=13, bold=True, color=NAVY, font=FONT_HEAD)
        add_textbox(slide, Emu(int(bx) + Inches(0.2)), Emu(int(by) + Inches(0.42)),
                    Emu(int(block_w) - Inches(0.2)), Inches(0.45),
                    desc, size=9.5, color=TEXT_MUTED)

    bars_top = Emu(int(top0) + 2 * int(block_h) + int(row_gap) + Inches(0.15))
    add_callout_bar(slide, MARGIN, bars_top, Inches(9.1), Inches(0.4), "Stable Control Plane",
                     "API Server, Scheduler, DAG Processor, Triggerer and PostgreSQL showed "
                     "limited resource increases.")
    add_callout_bar(slide, MARGIN, Emu(int(bars_top) + Inches(0.48)), Inches(9.1), Inches(0.42),
                     "", "Result: Airflow successfully orchestrated the pipeline within "
                     "configured worker and control-plane limits.",
                     bg_color=RGBColor(0xFC, 0xEC, 0xDA), text_color=TEXT_DARK, text_bold=True)


def add_governance_access_slide():
    slide = content_slide("Governance & Trust", "Access, Privacy & Auditability",
                           title_size=19)
    items = [
        ("Role-Based Access Control", "Enforced for every user across the datastore."),
        ("Tenant Row-Level Security", "Strict multi-tenant isolation at the row level."),
        ("Least-Privilege Writes", "Write access limited to orchestration services and DBAs."),
        ("PII Excluded by Design", "Sensitive personal information never reaches analytics."),
        ("Jurisdiction-Based Access", "Data access is scoped and restricted by jurisdiction."),
        ("Full Query Traceability", "Every query logged and attributable; runaway-query guardrails."),
    ]
    add_numbered_list(slide, items, MARGIN, CONTENT_TOP, Inches(5.9), Inches(3.85))
    add_icon_panel(slide, MSO_SHAPE.OCTAGON, NAVY, Inches(8.05), Inches(3.0), Inches(1.6))


def add_governance_ai_slide():
    slide = content_slide("Governance & Trust", "Responsible AI & Data Classification",
                           title_size=18)
    items = [
        ("Purpose Classification", "Every dataset tagged Operational, Policy Evaluation, or Research/AI."),
        ("AI Eligibility Tagging", "Datasets marked AI-Eligible, AI-Restricted, or AI-Prohibited."),
        ("Anonymization for AI", "Tokenization, salted hashing, and generalization before any AI use."),
        ("Human-in-the-Loop Approval", "Research/AI access needs governance committee sign-off, time-bound."),
        ("Versioned & Traceable", "Datasets like pt_silver_v3 carry batch-level lineage end to end."),
    ]
    add_numbered_list(slide, items, MARGIN, CONTENT_TOP, Inches(5.9), Inches(3.85))
    add_icon_panel(slide, MSO_SHAPE.SUN, ORANGE, Inches(8.05), Inches(3.0), Inches(1.6))


def add_schema_evolution_slide():
    slide = content_slide("Governance & Trust", "Schema Evolution, Handled Safely")
    left_w = Inches(4.35)
    left_x = MARGIN
    right_x = Emu(int(MARGIN) + int(left_w) + Inches(0.3))

    add_textbox(slide, left_x, CONTENT_TOP, left_w, Inches(0.4),
                "Raw Layer — Flexible by Design", size=13.5, bold=True, color=ORANGE_DARK,
                font=FONT_HEAD)
    add_dash_list(slide, [
        "Semi-structured JSON, no fixed column binding",
        "Upstream field changes never break ingestion",
        "Original payload preserved for replay",
    ], left_x, Emu(int(CONTENT_TOP) + Inches(0.45)), left_w, Inches(3.0), size=11)

    add_textbox(slide, right_x, CONTENT_TOP, left_w, Inches(0.4),
                "Silver & Gold — Versioned Contracts", size=13.5, bold=True, color=NAVY,
                font=FONT_HEAD)
    add_dash_list(slide, [
        "Any structural or KPI-logic change increments the version",
        "e.g. property_address_entity_v1 → v2",
        "Dependent marts version alongside (mart_defaulters_v1 → v2)",
        "Prior versions stay live through a controlled transition window",
    ], right_x, Emu(int(CONTENT_TOP) + Inches(0.45)), left_w, Inches(3.4), size=11)

    add_rect(slide, Emu(int(MARGIN) + int(left_w) + Inches(0.12)), CONTENT_TOP,
              Pt(1), Inches(3.6), RULE_GREY)


def add_value_props_slide():
    slide = content_slide("Why This Matters", "Three Reasons This Architecture Wins",
                           title_size=19)
    items = [
        ("Replayable & Auditable", "Raw data is immutable — silver can always be rebuilt "
         "from bronze if logic changes."),
        ("Fast, Reliable BI", "Gold marts are pre-aggregated — dashboards read finished "
         "numbers at constant latency."),
        ("Room for Analysts to Explore", "Silver is directly queryable — new questions "
         "don't need a pipeline change first."),
    ]
    add_numbered_list(slide, items, MARGIN, Inches(1.25), Inches(9.1), Inches(2.15))
    add_table(slide, ["Layer", "Primary Audience", "Typical Use"],
              [["Bronze", "Data engineers", "Debugging, replay, audit"],
               ["Silver", "Data / business analysts", "Ad hoc, exploratory queries"],
               ["Gold", "Business analysts, leadership", "Dashboards, reporting"]],
              left=MARGIN, top=Inches(3.65), width=Inches(9.1), row_h=0.32,
              col_widths=[1.9, 3.4, 3.8])


def add_roadmap_slide():
    slide = new_slide()
    set_background(slide, NAVY)
    add_textbox(slide, MARGIN, Inches(0.35), Inches(6), Inches(0.3),
                "ROADMAP", size=12, bold=True, color=ORANGE, font=FONT_HEAD)
    add_textbox(slide, MARGIN, Inches(0.68), Inches(9), Inches(0.55),
                "Where SakshamIQ Goes Next", size=24, bold=True, color=WHITE, font=FONT_HEAD)

    phases = [
        ("1", "Punjab Rollout", "Analytical datastore\nlive statewide", "March"),
        ("2", "Open Data Sharing", "Interoperable, open-\nformat data exchange", "TBD"),
        ("3", "Module Expansion", "Onboard more urban\nmodules onto the platform", "TBD"),
        ("4", "AI-Driven Capabilities", "Introduce AI/ML on\ntop of governed data", "TBD"),
        ("5", "Standardization", "Common data models\n& KPIs statewide", "TBD"),
    ]
    n = len(phases)
    top = Inches(2.55)
    node_d = Inches(0.55)
    gap = Inches(0.2)
    avail_w = Emu(int(CONTENT_RIGHT) - int(MARGIN))
    node_w = Emu((int(avail_w) - (n - 1) * int(gap)) // n)

    line_y = Emu(int(top) + int(node_d) // 2)
    add_rect(slide, MARGIN, line_y, Emu(int(avail_w)), Pt(1.5), NAVY_LIGHT)

    for i, (num, title, desc, timing) in enumerate(phases):
        cx = Emu(int(MARGIN) + i * (int(node_w) + int(gap)))
        node_left = Emu(int(cx) + (int(node_w) - int(node_d)) // 2)
        add_icon(slide, node_left, top, node_d, num, ORANGE, NAVY_DARK)
        add_textbox(slide, cx, Emu(int(top) + int(node_d) + Inches(0.18)), node_w, Inches(0.4),
                    title, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=FONT_HEAD)
        add_textbox(slide, cx, Emu(int(top) + int(node_d) + Inches(0.58)), node_w, Inches(0.75),
                    desc, size=8.5, color=RGBColor(0xC9, 0xD2, 0xEC), align=PP_ALIGN.CENTER)
        add_textbox(slide, cx, Emu(int(top) - Inches(0.35)), node_w, Inches(0.28),
                    timing, size=10, bold=True, color=ORANGE, align=PP_ALIGN.CENTER, font=FONT_HEAD)

    add_logo(slide, dark_bg=True)


def add_closing_slide():
    slide = new_slide()
    set_background(slide, WHITE)
    add_logo(slide)
    add_textbox(slide, Inches(0.8), Inches(2.0), Inches(8.4), Inches(0.8),
                "Thank You", size=36, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
                font=FONT_HEAD)
    add_rect(slide, Inches(4.55), Inches(2.85), Inches(0.9), Pt(2.5), ORANGE)
    add_textbox(slide, Inches(0.8), Inches(3.05), Inches(8.4), Inches(0.5),
                "SakshamIQ — Questions & Discussion", size=13, color=TEXT_MUTED,
                align=PP_ALIGN.CENTER)


# ------------------------------------------------------------------ build --

def build(output_path, template_path):
    global _PRS, _BLANK_LAYOUT
    _PRS = Presentation(template_path)
    delete_all_slides(_PRS)
    _PRS.slide_width = SLIDE_WIDTH
    _PRS.slide_height = SLIDE_HEIGHT
    _BLANK_LAYOUT = get_blank_layout(_PRS)

    add_title_slide()

    add_section_divider("Part One", "The Problem")
    add_problem_slide()
    add_why_punjab_slide()
    add_requirements_slide()
    add_why_clickhouse_slide()
    add_why_stack_slide()

    add_section_divider("Part Two", "The Architecture")
    add_principles_slide()
    add_diagram_slide()
    add_layer_detail_slide(
        "Bronze — Raw Events", ORANGE_DARK, MSO_SHAPE.CIRCULAR_ARROW,
        [
            ("Full, Untouched Event Log", "property_events_raw, demand_events_raw store "
             "every event exactly as received."),
            ("Append-Only", "Plain MergeTree — nothing is ever overwritten."),
            ("No Parsing at Ingest", "Upstream JSON is stored as-is; parsing happens later."),
            ("Preserved for Replay", "Every event ever delivered is kept, for replay and audit."),
        ],
    )
    add_layer_detail_slide(
        "Silver — Canonical Entities", NAVY, MSO_SHAPE.CUBE,
        [
            ("Populated Nightly", "Airflow parses bronze JSON and upserts entity tables."),
            ("Deduplicated", "ReplacingMergeTree, versioned on last_modified_time."),
            ("Full Entity Set", "Property, unit, owner, demand, payment, bill, assessment, audit."),
            ("Latest-State Queries", "Queried with FINAL whenever current state is required."),
        ],
    )
    add_layer_detail_slide(
        "Gold — Business Marts", ORANGE, MSO_SHAPE.LIGHTNING_BOLT,
        [
            ("Precomputed KPIs", "Built exclusively from canonical (silver) entities."),
            ("Broad Coverage", "Active properties, collections, defaulters, risk, assessments."),
            ("Dependency-Ordered Refresh", "One Refreshable Materialized View at a time."),
            ("Consistent Snapshots", "Dashboards never read a half-refreshed mart."),
        ],
    )
    add_billing_events_problem_slide()
    add_billing_code_changes_slide()
    add_billing_flow_slide()
    add_orchestration_slide()
    add_dag_bronze_to_silver_slide()
    add_dag_silver_to_gold_slide()

    add_section_divider("Part Three", "Validated at Scale")
    add_scale_slide()
    add_storage_slide()
    add_kpi_performance_slide()
    add_validated_airflow_slide()

    add_section_divider("Part Four", "Governance & Trust")
    add_governance_access_slide()
    add_governance_ai_slide()
    add_schema_evolution_slide()

    add_value_props_slide()
    add_roadmap_slide()
    add_closing_slide()

    _PRS.save(output_path)
    print(f"Wrote {output_path} ({len(_PRS.slides)} slides)")


def main():
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--output",
        default="docs/SakshamIQ.pptx",
        help="Output .pptx path (default: docs/SakshamIQ.pptx)",
    )
    parser.add_argument(
        "--template",
        default=os.path.expanduser("~/Downloads/Template.pptx"),
        help="Base eGov PowerPoint template to build on "
             "(default: ~/Downloads/Template.pptx)",
    )
    args = parser.parse_args()
    build(args.output, args.template)


if __name__ == "__main__":
    main()
